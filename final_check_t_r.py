from lxml import etree
import pandas as pd
import time
import dataframe_image as dfi
from tqdm import tqdm
from time import sleep
import matplotlib.pyplot as plt
from matplotlib import pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from datetime import date
import progressbar
from colorama import init
from termcolor import colored
import logging

# importing dictionaries needed for checking constant values
from dictionary_t_r import dictionary as dict


# function that checks the values in Excel and xml files
def check(excelfile_t, excelfile_r, xmlfile):
    # opening xml file based on xml file name
    tree = etree.parse(f'{xmlfile}.xml')
    root=tree.getroot()

    # opening excel files (type1ne and rental) based on excel file name
    excel_t = pd.read_excel(f'{excelfile_t}.xlsx', sheet_name = 'sheet1', header=3)
    excel_r = pd.read_excel(f'{excelfile_r}.xlsx', sheet_name = 'sheet2', header=3)

    # opening lists for correct and incorrect values
    log=[]
    matched=[]

    # defining tree path for ContractOffer and Program
    Contract = tree.xpath('.//Contract')

    # Exceles for loops
    # type1

    # Contract
    excel_t_contract = excel_t[['Contract', 'Customer', 'Author', 'Rep','Delegate',
                            'Name', 'Customer2','Subtype','Organization', 'Currency', 'Start Date','End Date']]
    excel_t_contract = excel_t_contract.drop_duplicates()

    # Program
    excel_t_program = excel_t[['Contract', 'Customer','Program','S_date','E_date','T Timeframe']]
    excel_t_program = excel_t_program.drop_duplicates()

    # CFC
    excel_t_cfc = excel_t[['Contract', 'Customer','Program', 'Min Value', 'Min Volume']]                       
    excel_t_cfc = excel_t_cfc.drop_duplicates()

    # Product
    excel_t_product = excel_t[['Contract', 'Customer','Program','Product','P Start Date','P End Date','Payment', 'Shipping', 
                        'Reason','Tax', 'Test', 'Quantity','Units', 'Price']]
    excel_t_product = excel_t_product.drop_duplicates()    

    # CFD
    excel_t_cfd = excel_t[['Contract', 'Customer','Program', 'Organization', 'Currency', 'Unit', 'Volume', 'Type', 'Comp']]  
    excel_t_cfd = excel_t_cfd.drop_duplicates()    
    
    # RENTAL

    # Program
    excel_r_program = excel_r[['Contract', 'Customer','Program','S_date','E_date','Strategy Type']]
    excel_r_program = excel_r_program.drop_duplicates()

    # Rental
    excel_r_rental = excel_r[['Contract', 'Customer','Program','Product','P Start Date','P End Date','Payment', 'Reason',  'Quantity', 'TOR Date',
                            '#', 'Warranty', 'Billing', 'Ownership', 'PO',
                            'Serial #', 'Price']]
    excel_r_rental = excel_r_rental.drop_duplicates()

    # CFD
    excel_r_cfd = excel_r[['Contract', 'Customer','Program','Unit', 'Type', 'P Type', 'Frequency', 'B Day', 'Mode']]  
    excel_r_cfd = excel_r_cfd.drop_duplicates()

    # CONTRACT information ('.//Contract') 
    for contract in Contract:

        # Loop for going through all the  excel rows
        for index in tqdm(excel_t_contract.index, desc = colored("Update on Excel loop:", "red"), colour="green"):
            
            # Contract - key that allows to find Contract in Excel file
            if contract.attrib['Contract'] == excel_t_contract['Contract'][index]:
                
                # Customer 
                try:
                    if pd.to_numeric(contract.attrib['Customer']) == excel_t_contract['Customer'][index]:
                        matched.append(['matched Customer', index, excel_t_contract['Customer'][index],'Contract'])
                    else:
                        log.append(['mismatched Customer', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', 'Contract', contract.attrib['Customer'],excel_t_contract['Customer'][index]])
                except:
                    log.append(['no attribute Customer', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', '', ''])
            
                # Owner
                try:
                    if contract.attrib['Owner']== excel_t_contract['Author'][index]:
                        matched.append(['matched Owner', index, excel_t_contract['Customer'][index],'Contract'])
                    else:
                        log.append(['mismatched Owner', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', 'Contract', contract.attrib['Owner'],excel_t_contract['Author'][index]])
                except:
                    log.append(['no attribute Owner', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', 'Contract', '', ''])
            
                # Author
                try:
                    if contract.attrib['Author']== excel_t_contract['Rep'][index]:
                        matched.append(['matched Author', index, excel_t_contract['Customer'][index],'Contract'])
                    else:
                        log.append(['mismatched Author', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1','Contract', contract.attrib['Author'],excel_t_contract['Rep'][index]])
                except:
                    log.append(['no attribute Author', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', '', ''])

                # Delegate
                try: 
                    if contract.attrib['Delegate']== excel_t_contract['Delegate'][index]:
                        matched.append(['matched Delegate', index, excel_t_contract['Customer'][index], 'Contract'])
                    else:
                        log.append(['mismatched Delegate', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', 'Contract', contract.attrib['Delegate'], excel_t_contract['Delegate'][index]])
                except:
                    log.append(['no attribute Delegate', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', 'Contract', 'OPTIONAK', 'CORRECT'])

                #  name
                try:
                    if contract.attrib['Name'] == excel_t_contract['Name'][index]:
                        matched.append(['matched  Name', index, excel_t_contract['Customer'][index], 'Contract'])
                    else:
                        log.append(['mismatched  Name', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1','Contract', contract.attrib['DocName'], excel_t_contract['Name'][index]])
                except:
                    log.append(['no attribute  name', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', 'Contract', '', ''])

                # Customer
                try:
                    if contract.attrib['Customer'] == excel_t_contract['Customer'][index]:
                        matched.append(['matched Customer', index, excel_t_contract['Customer'][index], 'Contract'])
                    else:
                        log.append(['mismatched Customer', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', 'Contract', contract.attrib['CustomerContractId'], excel_t_contract['Customer'][index]])
                except:
                    log.append(['no attribute Customer', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', 'Contract', '', ''])

                #  Type
                try:
                    if contract.attrib['DocType']== dict.doc_type['type1']:
                        matched.append(['matched  Type', index, excel_t_contract['Customer'][index], 'Contract'])
                    else:
                        log.append(['mismatched  Type', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', 'Contract', contract.attrib['DocType'],dict.doc_type['Purchase_Based']])
                except:
                    log.append(['no attribute  Type', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', '', ''])

                #  Subtype
                try:
                    for key in dict.subtype:
                        if key == excel_t_contract['Subtype'][index]:
                            if contract.attrib['SubType']== dict.subtype[key]:
                                matched.append(['matched  Subtype', index, excel_t_contract['Customer'][index], 'Contract'])
                            else:
                                log.append(['mismatched  Subtype', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1','Contract', contract.attrib['CtrtSubType'], dict.subtype[key]])
                except:
                    log.append(['no attribute  Subtype', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', 'Contract', '',''])

                # Contract Domain
                try:
                    if contract.attrib['ContractDomain']== dict.domain['Commercial']:
                        matched.append(['matched Contract Domain', index, excel_t_contract['Customer'][index], 'Contract'])
                    else:
                        log.append(['mismatched Contract Domain', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1','Contract', contract.attrib['ContractDomain'], dict.domain['Commercial']])
                except:
                    log.append(['no attribute Contract Domain', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', '', ''])

                # Brand
                try:
                    if contract.attrib['Brand']== dict.brand['Brand']:
                        matched.append(['matched Brand', index, excel_t_contract['Customer'][index], 'Contract'])
                    else:
                        log.append(['mismatched Brand', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', contract.attrib['GenericBrandType'], dict.brand['Brand']])
                except:
                    log.append(['no attribute Brand', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', '', ''])

                # Method 
                try:
                    for key in dict.distribution:
                        if contract.attrib['Type']== dict.distribution[key]:
                            matched.append(['matched Method', index, excel_t_contract['Customer'][index], 'Contract'])
                        else:
                            log.append(['mismatched Method', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1','Contract', contract.attrib['Type'], dict.distribution['Direct']])
                except:
                    log.append(['no attribute Method', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', '', ''])

                #  Status
                try:
                    if contract.attrib['DocStatus']== dict.doc_status['status1']:
                        matched.append(['matched  Status', index, excel_t_contract['Customer'][index], 'Contract'])
                    else:
                        log.append(['mismatched  Status', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', contract.attrib['DocStatus'], dict.doc_status['Draft']])
                except:
                    log.append(['no attribute  Status', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', '', ''])

                # Organization
                try:
                    if pd.to_numeric(contract.attrib['Org']) == excel_t_contract['Organization'][index]:
                        matched.append(['matched Organization', index, excel_t_contract['Customer'][index], 'Contract'])
                    else:
                        log.append(['mismatched Organization', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1','Contract', contract.attrib['OrgUnitName'], excel_t_contract['Organization'][index]])
                except:
                    log.append(['no attribute Organization', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', '', ''])

                # Time zone 
                try:
                    for key in dict.timezone:
                        if contract.attrib['OrgUnitName'] == key:
                            if contract.attrib['Timezone'] == dict.timezone[key]:
                                matched.append(['matched Timezone', index, excel_t_contract['Customer'][index], 'Contract'])
                            else:
                                log.append(['mismatched Timezone', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1','Contract', contract.attrib['EffectiveTimezoneID'], dict.timezone[key]])
                except:
                    log.append(['no attribute Time zone', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', 'Contract', '', ''])

                # Locale
                try:
                    for key in dict.locale:
                        if contract.attrib['Name'] == key:
                            if contract.attrib['Locale']== dict.locale[key]:
                                matched.append(['matched Locale', index, excel_t_contract['Customer'][index], 'Contract'])
                            else:
                                log.append(['mismatched Locale', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', 'Contract', contract.attrib['EffectiveTimezoneID'], dict.locale[key]])
                except:
                    log.append(['no attribute Locale', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', '', ''])

                # Currency
                try:
                    if contract.attrib['Currency'] == excel_t_contract['Currency'][index]:
                        matched.append(['matched Currency', index, excel_t_contract['Customer'][index], 'Contract'])
                    else:
                        log.append(['mismatched Currency', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1','Contract',contract.attrib['Currency'], excel_t_contract['Currency'][index]])
                except:
                    log.append(['no attribute Currency', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', 'Contract', '', ''])

                # Start Date
                try:
                    if pd.to_datetime(contract.attrib['Start Date']) == excel_t_contract['Start Date'][index]:
                        matched.append(['matched Start Date', index, excel_t_contract['Customer'][index], 'Contract'])
                    else:
                        log.append(['mismatched Start Date', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', pd.to_datetime(contract.attrib['StartDate']), excel_t_contract['Start Date'][index]])
                except:
                    log.append(['no attribute Start Date', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', '', ''])

                # End Date
                try:
                    if pd.to_datetime(contract.attrib['End Date']) == excel_t_contract['End Date'][index]:
                        matched.append(['matched End Date', index, excel_t_contract['Customer'][index], 'Contract'])
                    else:
                        log.append(['mismatched End Date', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', contract.attrib['EndDate'], excel_t_contract['End Date'][index]])
                except:
                    log.append(['no attribute End Date', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', '', ''])

                # Green
                try:
                    if contract.attrib['Green'] == dict.green['No']:
                        matched.append(['matched Green', index, excel_t_contract['Customer'][index], 'Contract'])
                    else:
                        log.append(['mismatched Green', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', contract.attrib['Green'], dict.Green['No']])
                except:
                    log.append(['no attribute Green', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','Contract', '', ''])

                
                # COMMITMENT information ('.//')
                 = contract.xpath('.// ')
                for commitment in :        
                    # Customer 
                    try:
                        if pd.to_numeric(commitment.attrib['Owner'][2:-5]) == excel_t_contract['Customer'][index]:
                            matched.append(['matched  Type', index, excel_t_contract['Customer'][index], ''])
                        else:
                            log.append(['mismatched  Type', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','',commitment.attrib['Owner'][2:-5], excel_t_contract['Customer'][index]])
                    except:
                        log.append(['no attribute  Type', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','','',''])
                    
                    #  Type
                    try:
                        if commitment.attrib['Type'] == dict.commitmenttype['Group']:
                            matched.append(['matched  Type', index, excel_t_contract['Customer'][index], ''])
                        else:
                            log.append(['mismatched  Type', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','', commitment.attrib['Type'], dict.commitmenttype['Group']])
                    except:
                        log.append(['no attribute  Type', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','','',''])
                    
                    # On Hold
                    try:
                        if commitment.attrib['On Hold'] == dict.onhold['No']:
                            matched.append(['matched On Hold', index, excel_t_contract['Customer'][index], ''])
                        else:
                            log.append(['mismatched On Hold', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','', commitment.attrib['OnHold'], dict.onhold['No']])
                    except:
                        log.append(['no attribute On Hold', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','','',''])
                    
                    # Index
                    try:
                        if commitment.attrib['Index'] == dict.index['index']:
                            matched.append(['matched Index', index, excel_t_contract['Customer'][index], ''])
                        else:
                            log.append(['mismatched Index', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1','', commitment.attrib['Index'], dict.Tindex['index']])
                    except:
                        log.append(['no attribute Index', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','','',''])
                    
                    # Start Date
                    try:
                        if pd.to_datetime(commitment.attrib['Start Date']) == excel_t_contract['Start Date'][index] and pd.to_datetime(commitment.attrib['Start Date']) == excel_r['Start Date'][index]:
                            matched.append(['matched Start Date', index, excel_t_contract['Customer'][index], ''])
                        else:
                            log.append(['mismatched  Start Date', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','',pd.to_datetime(commitment.attrib['Start Date']), excel_t_contract['Start Date'][index]])
                    except:
                        log.append(['no attribute Start Date', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','','',''])
                    
                    # End Date
                    try:
                        if pd.to_datetime(commitment.attrib['End Date']) == excel_t_contract['End Date'][index] and pd.to_datetime(commitment.attrib['End Date']) == excel_r['End Date'][index]:
                            matched.append(['matched  End Date', index, excel_t_contract['Customer'][index], ''])
                        else:
                            log.append(['mismatched  End Date', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', '',pd.to_datetime(commitment.attrib['End Date']), excel_t_contract['End Date'][index]])
                    except:
                        log.append(['no attribute End Date', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1', '','',''])
                    
                    #  Start Date
                    try:
                        if pd.to_datetime(commitment.attrib['StartDate']) == excel_t_contract['Start Date'][index] and pd.to_datetime(commitment.attrib['StartDate']) == excel_r['Start Date'][index]:
                            matched.append(['matched  Start Date', index, excel_t_contract['Customer'][index], ''])
                        else:
                            log.append(['mismatched  Start Date', index, excel_t_contract['Customer'][index],contract.attrib['Contract'],'type1','',pd.to_datetime(commitment.attrib['StartDate']), excel_t_contract['Start Date'][index]])
                    except:
                        log.append(['no attribute  Start Date', index, excel_t_contract['Customer'][index], contract.attrib['Contract'],'type1','','',''])
                    
    
        
        
        # PROGRAM information ('.//Program')
        Program = contract.xpath('.//Program')
       
        for program in Program:
            # CHECKING type1 PROGRAM ATTRIBUTES
            if program.attrib['ID'] == dict.programtype['type1']:
                # Loop for going through all the Contracts excel rows
                for index in excel_t_program.index:
                    sleep(0.0001)
                    # Contract - key that allows to find Contract Offer in Excel file
                    if contract.attrib['Contract'] == excel_t_program['Contract'][index]:
                        # Program - key that allows to find the Program in Excel file
                        if program.attrib['Name'] == excel_t_program['Program'][index]:
                          
                            # S_date
                            try:
                                if pd.to_datetime(program.attrib['Start Date']) == excel_t_program['S_date'][index]:
                                    matched.append(['matched S_date', index, excel_t_program['Customer'][index], 'Program'])
                                else:
                                    log.append(['mismatched S_date', index, excel_t_program['Customer'][index], contract.attrib['Contract'],'type1','Program', program.attrib['Star tDate'], excel_t_program['S_date'][index]])
                            except:
                                log.append(['no attribute S_date', index, excel_t_program['Customer'][index], contract.attrib['Contract'],'type1','Program','', ''])

                            # E_date
                            try:
                                if pd.to_datetime(program.attrib['End Date']) == excel_t_program['E_date'][index]:
                                    matched.append(['matched E_date', index, excel_t_program['Customer'][index], 'Program'])
                                else:
                                    log.append(['mismatched E_date', index, excel_t_program['Customer'][index], contract.attrib['Contract'],'type1','Program', program.attrib['End Date'], excel_t_program['E_date'][index]])
                            except:
                                log.append(['no attribute E_date', index, excel_t_program['Customer'][index], contract.attrib['Contract'],'type1','Program', '', ''])

                            # Business Segment 
                            try:
                                for key in dict.business_segment:
                                    if contract.attrib['Name'] == key:
                                        if program.attrib['Business Segment'] == dict.business_segment[key]:
                                            matched.append(['matched Business Segment', index, excel_t_program['Customer'][index], 'Program'])
                                        else:
                                            log.append(['mismatched Business Segment', index, excel_t_program['Customer'][index],contract.attrib['Contract'],'type1','Program',program.attrib['BusinessSegment'], dict.business_segment[key]])
                            except:
                                log.append(['no attribute Business Segment', index, excel_t_program['Customer'][index],contract.attrib['Contract'],'type1', 'Program', '', ''])

                            # T Timeframe 
                            try:
                                if program.attrib['T Timeframe'] == excel_t_program['T Timeframe'][index]:
                                    matched.append(['matched T Timeframe', index, excel_t_program['Customer'][index], 'Program'])
                                else:
                                    log.append(['mismatched T Timeframe', index, excel_t_program['Customer'][index],contract.attrib['Contract'],'type1', 'Program', program.attrib['T Timeframe'], excel_t_program['T Timeframe'][index]])
                            except:
                                log.append(['no attribute T Timeframe', index, excel_t_program['Customer'][index], contract.attrib['Contract'],'type1', 'Programm','', ''])

                            # Schedule Basis
                            try:
                                if program.attrib['Schedule Basis'] == dict.schedulebasis['Calendar Based']:
                                    matched.append(['matched Schedule Basis', index, excel_t_program['Customer'][index], 'Product'])
                                else:
                                    log.append(['mismatched Schedule Basis', index, excel_t_program['Customer'][index],contract.attrib['Contract'],'type1','Product', program.attrib['Schedule Basis'], dict.schedulebasis['Calendar Based']])
                            except:
                                log.append(['no attribute  Schedule Basis', index, excel_t_program['Customer'][index],contract.attrib['Contract'],'type1','Product', '',''])

                # CFC information ('.//CFComponents')
                Cfc = program.xpath('.//CFC')
                for cfc in Cfc:        
                    # Loop for going through all the Contracts excel rows
                    for index in excel_t_cfc.index:
                        # Contract - key that allows to find Contract Offer in Excel file
                        if contract.attrib['Contract'] == excel_t_cfc['Contract'][index]:
                            # Program - key that allows to find the Program in Excel file
                            if program.attrib['Name'] == excel_t_cfc['Program'][index]:
                          
                                # Sales Line From
                                try:
                                    if cfc.attrib['Sale Line'] == dict.salesline['All Transactions']:
                                        matched.append(['matched Sales Line From', index, excel_t_cfc['Customer'][index], 'CFC'])
                                    else:
                                        log.append(['mismatched Sales Line From', index, excel_t_cfc['Customer'][index], contract.attrib['Contract'],'type1','CFC', cfc.attrib['Sale Line'], dict.salesline['All Transactions']])
                                except:
                                    log.append(['no attribute Sales Line From', index, excel_t_cfc['Customer'][index], contract.attrib['Contract'],'type1','CFC', '', 'CORRECT'])

                                # Num of Ts
                                try:
                                    if cfc.attrib['Num Ts'] == dict.numTs['Num of Ts']:
                                        matched.append(['matched Num of Ts', index, excel_t_cfc['Customer'][index], 'CFC'])
                                    else:
                                        log.append(['mismatched Num of Ts', index, excel_t_cfc['Customer'][index], contract.attrib['Contract'],'type1','CFC',cfc.attrib['Num Ts'], dict.numTs['Num of Ts']])
                                except:
                                    log.append(['no attribute Num of Ts', index, excel_t_cfc['Customer'][index], contract.attrib['Contract'],'type1','CFC', '',''])

                                #  Min Value
                                try:
                                    if cfc.attrib['Name'] == dict.minvaluegpg['Min Value']:
                                        Cfd2 = cfc.xpath('.//CFD')
                                        try:
                                            for cfd in Cfd2:
                                                if pd.to_numeric(cfd.attrib['Value2']) == excel_t_cfc['Min Value'][index]:
                                                    matched.append(['matched  Min Value', index, excel_t_cfc['Customer'][index], 'CFC'])
                                                else:
                                                    log.append(['mismatched  Min Value', index, excel_t_cfc['Customer'][index], contract.attrib['Contract'],'type1','CFC',cfd.attrib['Value2'], excel_t_cfc['Min Value'][index]])
                                        except:
                                            log.append(['no value  Min Valuein Excel', index, excel_t_cfc['Customer'][index], contract.attrib['Contract'],'type1','CFC', '','CORRECT'])
                                except:
                                    log.append(['no attribute  Min Value', index, excel_t_cfc['Customer'][index], contract.attrib['Contract'],'type1','CFC', '','CORRECT'])

                                # VPG-1 Min Value
                                try:
                                    if cfc.attrib['Name'] == dict.minvaluevpg['Min Value']:
                                        Cfd2 = cfc.xpath('.//CFD')
                                        for cfd in Cfd2:
                                            if pd.to_numeric(cfd.attrib['Value2']) == excel_t_cfc['Min Volume'][index]:
                                                matched.append(['matched VPG-1 Min Value', index, excel_t_cfc['Customer'][index], 'CFC'])
                                            else:
                                                log.append(['mismatched VPG-1 Min Value', index, excel_t_cfc['Customer'][index], contract.attrib['Contract'],'type1','CFC',cfd.attrib['Value2'], excel_t_cfc['Min Volume'][index]])
                                except:
                                    log.append(['no attribute  Min Value', index, excel_t_cfc['Customer'][index], contract.attrib['Contract'],'type1','CFC', '',''])
          
                            
                # PRODUCT inormation ('.//ProductLI')
                Product = program.xpath('.//Product')
                for product in Product:
                    # Loop for going through all the Contracts excel rows
                    for index in excel_t_product.index:
                        # Contract - key that allows to find Contract Offer in Excel file
                        if contract.attrib['Contract'] == excel_t_product['Contract'][index]:
                            # Program - key that allows to find the Program in Excel file
                            if program.attrib['Name'] == excel_t_product['Program'][index]:
                                # Product - key that allows to find Product Number in Excel file
                                if excel_t_product['Product'][index].replace("-", "") == product.attrib['ProductNum'].replace("-", ""):
                                    
                                   # Direct
                                    try:
                                        if product.attrib['Direct'] == dict.distribution['Direct']:
                                            matched.append(['matched Direct', index, excel_t_product['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched Direct', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1', 'Product', product.attrib['Direct'], dict.distribution['Direct']])
                                    except:
                                        log.append(['no attribute Direct', index, excel_t_product['Customer'][index], contract.attrib['Contract'],'type1','Product', '',''])

                                    # Awarded Status
                                    try:
                                        if product.attrib['Awarded'] == dict.awarded['status1']:
                                            matched.append(['matched Awarded Status', index, excel_t_product['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched Awarded Status', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1', 'Product',product.attrib['Awarded'], dict.awarded['Primary']])
                                    except:
                                        log.append(['no attribute Awarded Status', index, excel_t_product['Customer'][index], contract.attrib['Contract'],'type1','Product', '',''])

                                    # P Start Date
                                    try:
                                        if pd.to_datetime(product.attrib['Start Date']) == excel_t_product['P Start Date'][index]:
                                            matched.append(['matched P Start Date', index, excel_t_product['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched P Start Date', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1','Product', pd.to_datetime(product.attrib['Start Date']), excel_t_product['P Start Date'][index]])
                                    except:
                                        log.append(['no attribute P Start Date', index, excel_t_product['Customer'][index], contract.attrib['Contract'],'type1','Product', '',''])

                                    # P End Date
                                    try:
                                        if pd.to_datetime(product.attrib['End Date']) == excel_t_product['P End Date'][index]:
                                            matched.append(['matched P End Date', index, excel_t_product['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched P End Date', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1', 'Product',pd.to_datetime(product.attrib['End Date']), excel_t_product['P End Date'][index]])
                                    except:
                                        log.append(['no attribute P End Date', index, excel_t_product['Customer'][index], contract.attrib['Contract'],'type1','Product', '',''])

                                    # Payment
                                    try:
                                        for key in dict.paymentterms:
                                            if key == product.attrib['Payment']:
                                                try:
                                                    if dict.paymentterms[key] == excel_t_product['Payment'][index]:
                                                        matched.append(['matched Payment', index, excel_t_product['Customer'][index], 'Product'])
                                                    else:
                                                        log.append(['mismatched Payment', index, excel_t_product['Customer'][index], contract.attrib['Contract'],'type1','Product',product.attrib['Payment'][5:8], excel_t_product['Payment'][index][:3]])
                                                except:
                                                    log.append(['no value Payment in Excel', index, excel_t_product['Customer'][index], contract.attrib['Contract'],'type1','Product', 'OPTIONAL',''])
                                    except: 
                                        log.append(['no attribute Payment', index, excel_t_product['Customer'][index], contract.attrib['Contract'],'type1','Product', 'OPTIONAL','CORRECT'])

                                    # Shipping
                                    try:
                                        if product.attrib['Shipping].replace('_', ' ') == excel_t_product['Shipping'][index].replace(':', '').upper():
                                            matched.append(['matched Shipping', index, excel_t_product['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched Shipping', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1','Product', program.attrib['Shipping'].replace('_', ' '), excel_t_product['Shipping'][index].replace(':', '').upper()])
                                    except:
                                        log.append(['no attribute Shipping', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1', 'Product', 'OPTIONAL','CORRECT'])

                                    # Reason
                                    try:
                                        if product.attrib['Reason'][:4] == excel_t_product['Reason'][index][:4]:
                                            matched.append(['matched Reason', index, excel_t_product['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched Reason', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1','Product', product.attrib['Reason'][:4], excel_t_product['Reason'][index][:4]])
                                    except:
                                        log.append(['no attribute Reason', index, excel_t_product['Customer'][index], contract.attrib['Contract'],'type1','Product', 'OPTIONAL','CORRECT'])

                                    # Tax
                                    try:
                                        if product.attrib['Tax'] == excel_t_product['Tax'][index].upper():
                                            matched.append(['matched Tax', index, excel_t_product['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched Tax', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1','Product', product.attrib['Tax'], excel_t_product['Tax'][index].upper()])
                                    except:
                                        log.append(['no attribute Tax', index, excel_t_product['Customer'][index], contract.attrib['Contract'],'type1','Product', 'OPTIONAL','CORRECT'])

                                    # Test
                                    try:
                                        if pd.to_datetime(product.attrib['Test']) == excel_t_product['Test'][index]:
                                            matched.append(['matched Test', index, excel_t_product['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched Test', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1','Product', pd.to_datetime(product.attrib['Test']),excel_t_product['Test'][index]])
                                    except:
                                        log.append(['no attribute  Test', index, excel_t_product['Customer'][index], contract.attrib['Contract'],'type1','Product','OPTIONAL','CORRECT'])

                                    # Quantity
                                    try:
                                        if pd.to_numeric(product.attrib['Quantity']) == excel_t_product['Quantity'][index]:
                                            matched.append(['matched Quantity', index, excel_t_product['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched Quantity', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1','Product', pd.to_numeric(product.attrib['Quantity']), excel_t_product['Quantity'][index]])
                                    except:
                                        log.append(['no attribute  Quantity', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1','Product', 'OPTIONAL','CORRECT'])
 
                                    # PRICE information ('.//type1Prices')
                                    Price = product.xpath('.//Price')
                                    for price in Price:
                               
                                        # Units
                                        try:
                                            if pd.to_numeric(price.attrib['Units']) == excel_t_product['Units'][index]:
                                                matched.append(['matched Units', index, excel_t_product['Customer'][index], 'Price'])
                                            else:
                                                log.append(['mismatched Units', index, excel_t_product['Customer'][index], contract.attrib['Contract'],'type1','Price', pd.to_numeric(price.attrib['Units']) == excel_t_product['Units'][index]])
                                        except:
                                            log.append(['no attribute Units', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1', 'Price', 'OPTIONAL',''])
                                        
                                        # Price 
                                        try:
                                            if pd.to_numeric(price.attrib['Price'][:-4]) == excel_t_product['Price'][index]:
                                                matched.append(['matched Price', index, excel_t_product['Customer'][index], 'Price'])
                                            else:
                                                log.append(['mismatched Price', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1', 'Price', pd.to_numeric(price.attrib['Price'][:-4]), excel_t_product['Price'][index]])
                                        except:
                                            log.append(['no attribute Price', index, excel_t_product['Customer'][index],contract.attrib['Contract'],'type1', 'Price', '',''])
                                        
                # CFD information ('.//CFD')
                Cfd = program.xpath('.//CFD')
                for cfd in Cfd:
                    # Loop for going through all the Contracts excel rows
                    for index in excel_t_cfd.index:
                        # Contract - key that allows to find Contract Offer in Excel file
                        if contract.attrib['Contract'] == excel_t_cfd['Contract'][index]:
                            # Program - key that allows to find the Program in Excel file
                            if program.attrib['Name'] == excel_t_cfd['Program'][index]:
                                # Value1 
                                try:
                                    if cfd.attrib['Value1']== dict.pricelist_value1['value1']:
                                        for key in dict.pricelist_value2:
                                            if key == str(excel_t_cfd['Organization'].apply(str)[index]+ '_' + excel_t_cfd['Currency'].apply(str)[index]):
                                                if cfd.attrib['Value2'] == dict.pricelist_value2[key]:
                                                    matched.append(['matched value1', index, excel_t_cfd['Customer'][index], 'CFD'])
                                                else:
                                                    log.append(['mismatched value1', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', cfd.attrib['Value2'], dict.pricelist_value2[key]])
                                except:
                                    log.append(['no attribute value1', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])

                                # Allow
                                try:
                                    if cfd.attrib['Value1'] == dict.allowexpired['Allow']:
                                        if cfd.attrib['Value2']== dict.allowexpired['Yes']:
                                            matched.append(['matched Allow', index, excel_t_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Allow', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1','CFD', cfd.attrib['Value1'], dict.allowexpired['Allow']])
                                except:
                                    log.append(['no attribute Allow', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])

                                # Unit od Pricing 
                                try:
                                    if cfd.attrib['Value1'] == dict.unitofpricing_value1['Unit']:
                                        for key in dict.unitofpricing_value2:
                                            if excel_t_cfd['Unit'][index] == key: 
                                                if cfd.attrib['Value2'] == dict.unitofpricing_value2[key]:
                                                    matched.append(['matched Unit', index, excel_t_cfd['Customer'][index], 'CFD'])
                                                else:
                                                    log.append(['mismatched Unit', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD',dict.unitofpricing_value2[key], excel_t_cfd['Unit'][index]])
                                except:
                                    log.append(['no attribute Unit od Pricing ', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])

                                # Adjust Type
                                try:
                                    if cfd.attrib['Value1'] == dict.adjusttype_t['Adjust Type']:
                                        if cfd.attrib['Value2']== dict.adjusttype_t['Fixed Amt']:
                                            matched.append(['matched Adjust Type', index, excel_t_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Adjust Type', index, excel_t_cfd['Customer'][index], contract.attrib['Contract'],'type1','CFD',cfd.attrib['Value2'], dict.adjusttype_t['Fixed Amt']])
                                except:
                                    log.append(['no attribute Adjust Type', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])

                                # Access Price
                                try:
                                    if cfd.attrib['Value1'] == dict.accessprice['Access Price']:
                                        if cfd.attrib['Value2']== dict.accessprice['Yes']:
                                            matched.append(['matched Access Price', index, excel_t_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Access Price', index, excel_t_cfd['Customer'][index], contract.attrib['Contract'],'type1','CFD',cfd.attrib['Value1'], dict.accessprice['Access Price']])
                                except:
                                    log.append(['no attribute Access Price', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])

                                # Price Program Subtype
                                try:
                                    if cfd.attrib['Value1'] == dict.priceprogramsubtype['Price Program Subtype']:
                                        if cfd.attrib['Value2']== dict.priceprogramsubtype['Yes']:
                                            matched.append(['matched Price Program Subtype', index, excel_t_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Price Program Subtype', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD',cfd.attrib['Value1'], dict.priceprogramsubtype['Price Program Subtype']])
                                except:
                                    log.append(['no attribute Price Program Subtype', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])
                                
                                # Volume
                                try:
                                    if cfd.attrib['Value1'] == dict.volumetimeframe['Volume']:
                                        if cfd.attrib['Value2']== excel_t_cfd['Volume'][index]:
                                            matched.append(['matched Volume Timframe', index, excel_t_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Volume', index, excel_t_cfd['Customer'][index], contract.attrib['Contract'],'type1','CFD',cfd.attrib['Value1'], dict.volumetimeframe['Volume']])
                                except:
                                    log.append(['no attribute Volume', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])

                                # Type
                                try:
                                    if cfd.attrib['Value1'] == dict.producttypepriced_value1['Type']:
                                        for key in dict.producttypepriced_value2:
                                            if key == excel_t_cfd['Type'][index]:
                                                if cfd.attrib['Value2'] == dict.producttypepriced_value2[key]:
                                                    matched.append(['matched Type', index, excel_t_cfd['Customer'][index], 'CFD'])
                                                else:
                                                    log.append(['mismatched Type', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1','CFD', cfd.attrib['Value2'], excel_t_cfd['Type'][index]])
                                except:
                                    log.append(['no attribute Type', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])

                                # PBID
                                try:
                                    if cfd.attrib['Value1'] == dict.pbid['value1']:
                                        if cfd.attrib['Value2']== dict.pbid['No']:
                                            matched.append(['matched PBID', index, excel_t_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched PBID', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD',cfd.attrib['Value2'], dict.pbid['No']])
                                except:
                                    log.append(['no attribute PBID', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])

                                # T Rule
                                try:
                                    if cfd.attrib['Value1'] == dict.Trule['T Rule']:
                                        if cfd.attrib['Value2']== dict.Trule['Min T Attained']:
                                            matched.append(['matched T Rule', index, excel_t_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched T Rule', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD',cfd.attrib['Value1'], dict.Trule['T Rule']])
                                except:
                                    log.append(['no attribute T Rule', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])

                                # Comp
                                try:
                                    if cfd.attrib['Value1'] == dict.compliance_threshold['Comp']:
                                        try:
                                            if pd.to_numeric(cfd.attrib['Value2']) == pd.to_numeric(excel_t_cfd['Comp'][index]):
                                                matched.append(['matched Comp', index, excel_t_cfd['Customer'][index], 'CFD'])
                                            else:
                                                log.append(['mismatched Comp', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1','CFD', pd.to_numeric(cfd.attrib['Value2']), pd.to_numeric(excel_t_cfd['Comp'][index])])
                                        except:
                                            log.append(['no value Comp in excel_t_cfd', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', 'OPTIONAL', 'CORRECT'])
                                except:
                                    log.append(['no attribute Comp', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', 'OPTIONAL', 'CORRECT'])

                                # Compliance 
                                try:
                                    if cfd.attrib['Value1'] == dict.compliance_tracking_value1['Compliance ']:
                                        try:
                                            for key in dict.compliance_tracking_value2:
                                                if excel_t_cfd['Compliance Y/N'][index] == key:
                                                    if cfd.attrib['Value2'] == dict.compliance_tracking_value2[key]:
                                                        matched.append(['matched Compliance ', index, excel_t_cfd['Customer'][index], 'CFD'])
                                                    else:
                                                        log.append(['mismatched Compliance ', index, excel_t_cfd['Customer'][index], contract.attrib['Contract'],'type1','CFD',cfd.attrib['Value2'], dict.compliance_tracking_value2['Yes']])
                                        except:
                                            log.append(['no value Compliance  in Excel', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', 'OPTIONAL', ''])
                                except:
                                    log.append(['no attribute Compliance ', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', 'OPTIONAL', ''])
                                
                                # Market Basket Type
                                try:
                                    if cfd.attrib['Value1'] == dict.marketbasket['Market Basket Type']:
                                        if cfd.attrib['Value2']== dict.marketbasket['Dynamic']:
                                            matched.append(['matched Market Basket Type', index, excel_t_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Market Basket Type', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD',cfd.attrib['Value1'], dict.marketbasket['Market Basket Type']])
                                except:
                                    log.append(['no attribute Market Basket Type', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])

                                # Increase Rule 
                                try:
                                    if cfd.attrib['Value1'] == dict.increaserule['Increase Rule']:
                                        if cfd.attrib['Value2']== dict.increaserule['Value2']:
                                            matched.append(['matched Increase Rule', index, excel_t_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Increase Rule', index, excel_t_cfd['Customer'][index], contract.attrib['Contract'],'type1','CFD',cfd.attrib['Value1'], dict.increaserule['Increase Rule']])
                                except:
                                    log.append(['no attribute Increase Rule ', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])

                                # Increase Cap 
                                try:
                                    if cfd.attrib['Value1'] == dict.increasecap['Increase Cap']:
                                        if cfd.attrib['Value2']== dict.increasecap['Value2']:
                                            matched.append(['matched Increase Cap', index, excel_t_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Increase Cap', index, excel_t_cfd['Customer'][index], contract.attrib['Contract'],'type1','CFD',cfd.attrib['Value1'], dict.increasecap['Increase Cap']])
                                except:
                                    log.append(['no attribute Increase Cap ', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])

                                # Index Type 
                                try:
                                    if cfd.attrib['Value1'] == dict.indextype['Index Type']:
                                        if cfd.attrib['Value2']== dict.indextype['Value2']:
                                            matched.append(['matched Index Type', index, excel_t_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Index Type', index, excel_t_cfd['Customer'][index], contract.attrib['Contract'],'type1','CFD',cfd.attrib['Value1'], dict.indextype['Index Type']])
                                except:
                                    log.append(['no attribute Index Type', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])  

                                # Send T 
                                try:
                                    if cfd.attrib['Value1'] == dict.sendT1_value1['Send T1 Eligibility On Bid Award']:
                                        if cfd.attrib['Value2']== dict.sendT1_value2['True']:
                                            matched.append(['matched Send T ', index, excel_t_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Send T ', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD',cfd.attrib['Value2'], dict.sendT1_value2['True']])
                                except:
                                    log.append(['no attribute Send T ', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])
                                
                                # Measurement OFFset
                                try:
                                    if cfd.attrib['Value1'] == dict.measurement_value1['Measurement Offset']:
                                        if cfd.attrib['Value2']== dict.measurement_value2['Current Period']:
                                            matched.append(['matched Measurement OFFset', index, excel_t_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Measurement OFFset', index, excel_t_cfd['Customer'][index], contract.attrib['Contract'],'type1','CFD',cfd.attrib['Value2'], dict.measurement_value2['Current Period']])
                                except:
                                    log.append(['no attribute Measurement OFFset', index, excel_t_cfd['Customer'][index],contract.attrib['Contract'],'type1', 'CFD', '', ''])

            # CHECKING RENTAL AND SERVICE PROGRAM ATTRIBUTES
            elif program.attrib['SourceStrategyID'] == dict.programtype['rental'] or program.attrib['SourceStrategyID'] == dict.programtype['service'] or program.attrib['SourceStrategyID'] == dict.programtype['subscription']:
        
                for index in excel_r_program.index:
                    sleep(0.0001)
                    # Contract - key that allows to find Contract Offer in Excel file
                    if contract.attrib['Contract'] == excel_r_program['Contract'][index]:
                    #  Program - key that allows to find the Program in Excel file
                        if program.attrib['Name'] == excel_r_program['Program'][index]:
                            # S_date
                            try:
                                if pd.to_datetime(program.attrib['ContractStartDate']) == excel_r_program['S_date'][index]:
                                    matched.append(['matched S_date', index, excel_r_program['Customer'][index], 'Program'])
                                else:
                                    log.append(['mismatched S_date', index, excel_r_program['Customer'][index], contract.attrib['Contract'], 'Rental','Program', program.attrib['ContractStartDate'], excel_r_program['S_date'][index]])
                            except:
                                log.append(['no attribute S_date', index, excel_r_program['Customer'][index], contract.attrib['Contract'], 'Rental','Program','', ''])
                            
                            # E_date
                            try:
                                if pd.to_datetime(program.attrib['EndDate']) == excel_r_program['E_date'][index]:
                                    matched.append(['matched E_date', index, excel_r_program['Customer'][index], 'Program'])
                                else:
                                    log.append(['mismatched E_date', index, excel_r_program['Customer'][index], contract.attrib['Contract'], 'Rental','Program', program.attrib['EndDate'], excel_r_program['E_date'][index]])
                            except:
                                log.append(['no attribute E_date', index, excel_r_program['Customer'][index], contract.attrib['Contract'], 'Rental','Program', '', ''])
                            
                            # Business Segment 
                            try:
                                for key in dict.business_segment:
                                    if contract.attrib['OrgUnitName'] == key:
                                        if program.attrib['BusinessSegment'] == dict.business_segment[key]:
                                            matched.append(['matched Business Segment', index, excel_r_program['Customer'][index], 'Program'])
                                        else:
                                            log.append(['mismatched Business Segment', index, excel_r_program['Customer'][index],contract.attrib['Contract'], 'Rental','Program',program.attrib['BusinessSegment'], dict.business_segment[key]])
                            except:
                                log.append(['no attribute Business Segment', index, excel_r_program['Customer'][index],contract.attrib['Contract'], 'Rental', 'Program', '', ''])
                            
                            # Strategy Type
                            try:
                                if program.attrib['SourceStrategyID'][6:-15] == excel_r_program['Strategy Type'][index]:
                                    matched.append(['matched Strategy Type', index, excel_r_program['Customer'][index], 'Program'])
                                else:
                                    log.append(['mismatched Strategy Type', index, excel_r_program['Customer'][index],contract.attrib['Contract'], 'Rental', 'Program', program.attrib['SourceStrategyID'][6:-15], excel_r_program['Strategy Type'][index]])
                            except:
                                log.append(['no attribute Strategy Type', index, excel_r_program['Customer'][index], contract.attrib['Contract'], 'Rental', 'Programm','OPTIONAL', 'CORRECT'])

                                    
                            # CFC information ('.//CFC')
                            Cfc = program.xpath('.//CFC')
                            for cfc in Cfc:        
                    
                                # Sales Line From
                                try:
                                    if cfc.attrib['SaleLineType'] == dict.salesline['All Transactions']:
                                        matched.append(['matched Sales Line From', index, excel_r_program['Customer'][index], 'CFC'])
                                    else:
                                        log.append(['mismatched Sales Line From', index, excel_r_program['Customer'][index], contract.attrib['Contract'], 'Rental','CFC', cfc.attrib['SaleLineType'], dict.salesline['All Transactions']])
                                except:
                                    log.append(['no attribute Sales Line From', index, excel_r_program['Customer'][index], contract.attrib['Contract'], 'Rental','CFC', '',''])
                        
                                # Num of Ts
                                try:
                                    if cfc.attrib['NumTs'] == dict.numTs['Num of Ts']:
                                        matched.append(['matched Num of Ts', index, excel_r_program['Customer'][index], 'CFC'])
                                    else:
                                        log.append(['mismatched Num of Ts', index, excel_r_program['Customer'][index], contract.attrib['Contract'], 'Rental','CFC',cfc.attrib['NumTs'], dict.numTs['Num of Ts']])
                                except:
                                    log.append(['no attribute Num of Ts', index, excel_r_program['Customer'][index], contract.attrib['Contract'], 'Rental','CFC', '',''])
                        
                        
                # RENTAL inormation ('.//Rental')
                Rental = program.xpath('.//Rental')
                for rental in Rental:
                    for index in excel_r_rental.index:
                        # Contract - key that allows to find Contract Offer in Excel file
                        if contract.attrib['Contract'] == excel_r_rental['Contract'][index]:
                            #  Program - key that allows to find the Program in Excel file
                            if program.attrib['Name'] == excel_r_rental['Program'][index]:
                                # Product - key that allows to find Product Number in excel_r_rental file
                                if excel_r_rental['Product'][index].replace("-", "") == rental.attrib['ProductNum'].replace("-", ""):
                                    # Direct
                                    try:
                                        if rental.attrib['Direct'] == dict.distribution['Direct']:
                                            matched.append(['matched Direct', index, excel_r_rental['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched Direct', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental', 'Product', rental.attrib['Direct'], dict.distribution['Direct']])
                                    except:
                                        log.append(['no attribute Direct', index, excel_r_rental['Customer'][index], contract.attrib['Contract'], 'Rental','Product', '',''])

                                    # Awarded Status
                                    try:
                                        if rental.attrib['Awarded'] == dict.awarded['Primary']:
                                            matched.append(['matched Awarded Status', index, excel_r_rental['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched Awarded Status', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental', 'Product',rental.attrib['Awarded'], dict.awarded['Primary']])
                                    except:
                                        log.append(['no attribute Awarded Status', index, excel_r_rental['Customer'][index], contract.attrib['Contract'], 'Rental','Product', '',''])

                                    # P Start Date
                                    try:
                                        if pd.to_datetime(rental.attrib['StartEffectiveDate']) == excel_r_rental['P Start Date'][index]:
                                            matched.append(['matched P Start Date', index, excel_r_rental['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched P Start Date', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental','Product', pd.to_datetime(rental.attrib['StartEffectiveDate']), excel_r_rental['P Start Date'][index]])
                                    except:
                                        log.append(['no attribute P Start Date', index, excel_r_rental['Customer'][index], contract.attrib['Contract'], 'Rental','Product', '',''])

                                    # P End Date
                                    try:
                                        if pd.to_datetime(rental.attrib['EndEffectiveDate']) == excel_r_rental['P End Date'][index]:
                                            matched.append(['matched P End Date', index, excel_r_rental['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched P End Date', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental', 'Product',pd.to_datetime(rental.attrib['EndEffectiveDate']), excel_r_rental['P End Date'][index]])
                                    except:
                                        log.append(['no attribute P End Date', index, excel_r_rental['Customer'][index], contract.attrib['Contract'], 'Rental','Product', '',''])

                                    # Payment
                                    try:
                                        for key in dict.paymentterms:
                                            if key == rental.attrib['PaymentTerm']:
                                                try:
                                                    if dict.paymentterms[key] == excel_r_rental['Payment'][index]:
                                                        matched.append(['matched Payment', index, excel_r_rental['Customer'][index], 'Product'])
                                                    else:
                                                        log.append(['mismatched Payment', index, excel_r_rental['Customer'][index], contract.attrib['Contract'], 'Rental','Product',rental.attrib['PaymentTerm'][5:8], excel_r_rental['Payment'][index][:3]])
                                                except:
                                                    log.append(['no value Payment in Excel', index, excel_r_rental['Customer'][index], contract.attrib['Contract'], 'Rental','Product', 'OPTIONAL',''])
                                    except:
                                        log.append(['no attribute  Payment', index, excel_r_rental['Customer'][index], contract.attrib['Contract'], 'Rental','Product', 'OPTIONAL','CORRECT'])

            
                                    # Reason
                                    try:
                                        if rental.attrib['ReasonCode'][:4] == excel_r_rental['Reason'][index][:4]:
                                            matched.append(['matched Reason', index, excel_r_rental['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched Reason', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental','Product', rental.attrib['ReasonCode'][:4], excel_r_rental['Reason'][index][:4]])
                                    except:
                                        log.append(['no attribute Reason', index, excel_r_rental['Customer'][index], contract.attrib['Contract'], 'Rental','Product', 'OPTIONAL','CORRECT'])

                    
                                    # Test
                                    try:
                                        if pd.to_datetime(rental.attrib['TestOfRecordDate']) == excel_r_rental['TOR Date'][index]:
                                            matched.append(['matched Test', index, excel_r_rental['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched Test', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental','Product', pd.to_datetime(rental.attrib['TestOfRecordDate']),excel_r_rental['Test'][index]])
                                    except:
                                        log.append(['no attribute  Test', index, excel_r_rental['Customer'][index], contract.attrib['Contract'], 'Rental','Product', 'OPTIONAL','CORRECT'])

                                    # Quantity
                                    try:
                                        if pd.to_numeric(rental.attrib['Quantity']) == excel_r_rental['Quantity'][index]:
                                            matched.append(['matched Quantity', index, excel_r_rental['Customer'][index], 'Product'])
                                        else:
                                            log.append(['mismatched TQuantity', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental','Product', pd.to_numeric(rental.attrib['Quantity']), excel_r_rental['Quantity'][index]])
                                    except:
                                        log.append(['no attribute  Quantity', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental','Product', 'OPTIONAL',''])

                                    # #
                                    try:
                                        if rental.attrib['CachedInstrumentProductNum'].replace('-', '') == excel_r_rental['#'][index].replace('-', ''):
                                            matched.append(['matched Instrument', index, excel_r_rental['Customer'][index], 'Rental'])
                                        else:
                                            log.append(['mismatched Instrument', index, excel_r_rental['Customer'][index], contract.attrib['Contract'],'Rental', 'Rental', rental.attrib['CachedInstrumentProductNum'].replace('-', ''), excel_r_rental['#'][index].replace('-', '')])
                                    except:
                                            log.append(['no attribute Instrument', index, excel_r_rental['Customer'][index], contract.attrib['Contract'],'Rental', 'Rental','',''])
                        
                                    # Warranty
                                    try:
                                        if pd.to_numeric(rental.attrib['Warranty']) == excel_r_rental['Warranty'][index]:
                                            matched.append(['matched Warranty', index, excel_r_rental['Customer'][index], 'Rental'])
                                        else:
                                            log.append(['mismatched Warranty', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental', 'Rental', rental.attrib['CachedInstrumentProductNum'].replace('-', ''), excel_r_rental['#'][index].replace('-', '')])
                                    except:
                                        log.append(['no attribute Warranty', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental', 'Rental','OPTIONAL','CORRECT'])
                                    
                                    # Billing
                                    try:
                                        if pd.to_datetime(rental.attrib['BillingStartDate']) == excel_r_rental['Billing'][index]:
                                            matched.append(['matched Billing', index, excel_r_rental['Customer'][index], 'Rental'])
                                        else:
                                            log.append(['mismatched Billing', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental', 'Rental', rental.attrib['BillingStartDate'], excel_r_rental['Billing'][index]])
                                    except:
                                        log.append(['no attribute Billing', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental','Rental', '', ''])

                                    # Ownership
                                    try:
                                        for key in dict.ownership:
                                            if dict.ownership[key] == rental.attrib['Ownership']:
                                                try:
                                                    if key  == excel_r_rental['Ownership'][index]:
                                                        matched.append(['matched Ownership', index, excel_r_rental['Customer'][index], 'Rental'])
                                                    else:
                                                        log.append(['no value Ownership in Excel', index, excel_r_rental['Customer'][index],contract.attrib['Contract'],'Rental', 'Rental', rental.attrib['Ownership'], excel_r_rental['Ownership'][index]])
                                                except:
                                                    log.append(['no attribute Ownership', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental','Rental', 'OPTIONAL', ''])
                                    except:
                                        log.append(['no attribute Ownership', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental','Rental', 'OPTIONAL', 'CORRECT'])
                        
                                    # PO
                                    try:
                                        if rental.attrib['PONum'] == excel_r_rental['PO'][index]:
                                            matched.append(['matched PO', index, excel_r_rental['Customer'][index], 'Rental'])
                                        else:
                                            log.append(['mismatched PO', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental', 'Rental', rental.attrib['PONum'], excel_r_rental['PO'][index]])
                                    except:
                                        log.append(['no attribute PO', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental','Rental', 'OPTIONAL', 'CORRECT'])

                                    # INCLUDED information ('.//IncludedLIAttr')
                                    Included = rental.xpath('.//IncludedLIAttr')
                                    for included in Included:
                                    # Serial #
                                        try:
                                            if included.attrib['Value'] == excel_r_rental['Serial #'][index]:
                                                matched.append(['matched Serial', index, excel_r_rental['Customer'][index], 'Included'])
                                            else:
                                                log.append(['mismatched Serial', index, excel_r_rental['Customer'][index], contract.attrib['Contract'],'Rental', 'Included', included.attrib['Value'], excel_r_rental['Serial #'][index]])
                                        except:
                                            log.append(['no attribute Serial', index, excel_r_rental['Customer'][index], contract.attrib['Contract'],'Rental', 'Included','OPTIONAL',''])
            
                                    # PRICE information ('.//type1Prices')
                                    Price = rental.xpath('.//LIPrices')
                                    for price in set(Price):
                                        # Price 
                                        try:
                                            if pd.to_numeric(price.attrib['PriceHigh'][:-4]) == excel_r_rental['Price'][index]:
                                                matched.append(['matched Price', index, excel_r_rental['Customer'][index], 'Price'])
                                            else:
                                                log.append(['mismatched Price', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental', 'Price', pd.to_numeric(price.attrib['PriceHigh'][:-4]), excel_r_rental['Price'][index]])
                                        except:
                                            log.append(['no attribute Price', index, excel_r_rental['Customer'][index],contract.attrib['Contract'], 'Rental', 'Price', '',''])
                                        
                # CFD information ('.//CFD')
                Cfd = program.xpath('.//CFD')
                for cfd in Cfd:
                    for index in excel_r_cfd.index:
                        # Contract - key that allows to find Contract Offer in Excel file
                        if contract.attrib['Contract'] == excel_r_cfd['Contract'][index]:
                            #  Program - key that allows to find the Program in Excel file
                            if program.attrib['Name'] == excel_r_cfd['Program'][index]:
                                # Unit od Pricing 
                                try:
                                    if cfd.attrib['Value1'] == dict.unitofpricing_value1['Unit']:
                                        for key in dict.unitofpricing_value2:
                                            if excel_r_cfd['Unit'][index] == key: 
                                                if cfd.attrib['Value2'] == dict.unitofpricing_value2[key]:
                                                    matched.append(['matched Unit', index, excel_r_cfd['Customer'][index], 'CFD'])
                                                else:
                                                    log.append(['mismatched Unit', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD',cfd.attrib['Value2'], dict.unitofpricing_value2[key]])
                                except:
                                    log.append(['no attribute Unit od Pricing ', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD', '', ''])

                                # Adjust Type
                                try:
                                    if cfd.attrib['Value1'] == dict.adjusttype_r['Adjust Type']:
                                        if cfd.attrib['Value2']== dict.adjusttype_r['Fixed Amt']:
                                            matched.append(['matched Adjust Type', index, excel_r_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Adjust Type', index, excel_r_cfd['Customer'][index], contract.attrib['Contract'], 'Rental','CFD',cfd.attrib['Value1'], dict.adjusttype_r['Adjust Type']])
                                except:
                                    log.append(['no attribute Adjust Type', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD', '', ''])
                        

                                # Type
                                try:
                                    if cfd.attrib['Value1'] == dict.producttypepriced_value1['Type']:
                                        for key in dict.producttypepriced_value2:
                                            if key == excel_r_cfd['Type'][index]:
                                                if cfd.attrib['Value2'] == dict.producttypepriced_value2[key]:
                                                    matched.append(['matched Type', index, excel_r_cfd['Customer'][index], 'CFD'])
                                                else:
                                                    log.append(['mismatched Type', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental','CFD', cfd.attrib['Value2'], excel_r_cfd['Type'][index]])
                                except:
                                    log.append(['no attribute Type', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD', '', ''])

                                # Increase Rule 
                                try:
                                    if cfd.attrib['Value1'] == dict.increaserule['Increase Rule']:
                                        if cfd.attrib['Value2']== dict.increaserule['Value2']:
                                            matched.append(['matched Increase Rule', index, excel_r_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Increase Rule', index, excel_r_cfd['Customer'][index], contract.attrib['Contract'], 'Rental','CFD',cfd.attrib['Value1'], dict.increaserule['Increase Rule']])
                                except:
                                    log.append(['no attribute Increase Rule ', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD', '', ''])

                                # Increase Cap 
                                try:
                                    if cfd.attrib['Value1'] == dict.increasecap['Increase Cap']:
                                        if cfd.attrib['Value2']== dict.increasecap['Value2']:
                                            matched.append(['matched Increase Cap', index, excel_r_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Increase Cap', index, excel_r_cfd['Customer'][index], contract.attrib['Contract'], 'Rental','CFD',cfd.attrib['Value1'], dict.increasecap['Increase Cap']])
                                except:
                                    log.append(['no attribute Increase Cap ', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD', '', ''])
                        
                        
                                # Send T 
                                try:
                                    if cfd.attrib['Value1'] == dict.sendT1_value1['Send T1 Eligibility On Bid Award']:
                                        if cfd.attrib['Value2']== dict.sendT1_value2['True']:
                                            matched.append(['matched Send T ', index, excel_r_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Send T ', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD',cfd.attrib['Value2'], dict.sendT1_value2['True']])
                                except:
                                    log.append(['no attribute Send T ', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD', '', ''])
                                
                                # P Type
                                try:
                                    if cfd.attrib['Value1'] == dict.pricingtype['P Type']:
                                        if cfd.attrib['Value2'].replace('_', ' ') == excel_r_cfd['P Type'][index].upper():
                                            matched.append(['matched P Type', index, excel_r_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched P Type', index, excel_r_cfd['Customer'][index], contract.attrib['Contract'], 'Rental', 'CFD',cfd.attrib['Value2'].replace('_', ' '), excel_r_cfd['P Type'][index].upper()])
                                except:
                                    log.append(['no attribute P Type', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD', '', ''])
                                
                                # Frequency
                                try:
                                    if cfd.attrib['Value1'] == dict.invoice['Frequency']:
                                        if cfd.attrib['Value2'] == excel_r_cfd['Frequency'][index][:3].upper():
                                            matched.append(['matched Frequency', index, excel_r_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched Frequency', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD',cfd.attrib['Value2'], excel_r_cfd['Frequency'][index][:3].upper()])
                                except:
                                    log.append(['no attribute Frequency', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD', '', ''])
                                
                                # B Day
                                try:
                                    if cfd.attrib['Value1'] == dict.billingday['B Day']:
                                        if pd.to_numeric(cfd.attrib['Value2']) == excel_r_cfd['B Day'][index]:
                                            matched.append(['matched B Day', index, excel_r_cfd['Customer'][index], 'CFD'])
                                        else:
                                            log.append(['mismatched B Day', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD',pd.to_numeric(cfd.attrib['Value2']), excel_r_cfd['B Day'][index]])
                                except:
                                    log.append(['no attribute B Day', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD', '', ''])
                                
                                # B Day Mode
                                try:
                                    if cfd.attrib['Value1'] == dict.billingdaymode_value1['B Day Mode']:
                                        for key in dict.billingdaymode_value2:
                                            try:
                                                if key == excel_r_cfd['Mode'][index]:
                                                    if cfd.attrib['Value2'] == dict.billingdaymode_value2[key]:
                                                        matched.append(['matched B Day Mode', index, excel_r_cfd['Customer'][index], 'CFD'])
                                                    else:
                                                        log.append(['mismatched B Day Mode', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD',cfd.attrib['Value2'], excel_r_cfd['Mode'][index]])
                                            except:
                                                log.append(['no value B Day Mode in Excel', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD', '',''])
                                except:
                                    log.append(['no attribute B Day Mode', index, excel_r_cfd['Customer'][index],contract.attrib['Contract'], 'Rental', 'CFD', '',''])


                    
    # Current time for saving the log and correct files
    timestr = time.strftime("%Y%m%d-%H%M%S")

    # Saving log and correct lists as dataframes and exporting them to csv files
    matched = pd.DataFrame((matched), columns =['Comment', 'Index in excel', 'Customer number', 'Area'])
    matched.to_csv(f'matched {timestr}.csv') 
    log = pd.DataFrame((log), columns =['Comment', 'Index in excel', 'Customer number', 'Alternate ID', 'Program Type', 'Area', 'Value from xml file', 'Value from Excel file/Constant value'])
    # Deleting correct optionals from log
    log = log.loc[log['Value from Excel file/Constant value'] != 'CORRECT']
    log.to_csv(f'{xmlfile[35:]}-log {timestr}.csv') 

    # Exporting log to png file as mytable (needed for powerpoint presentation)- grouped by the error
    log_grouped = log.groupby(['Comment', 'Customer number', 'Alternate ID', 'Program Type', 'Area', 'Value from xml file', 'Value from Excel file/Constant value']).count()
    log_grouped = log_grouped.rename(columns= {'Index in excel' : 'Sum of mismatches'})
    log_grouped.to_csv(f'{xmlfile[35:]}-log grouped.csv') 
    dfi.export(log_grouped,f"mytable1-{xmlfile[35:]}.png")

    
    # Counting number of correct and incorrect attributes and saving as dataframe
    rows = [['Number of matched attributes',matched.shape[0]], ['Number of mismatched attributes', log.shape[0]]]
    statistic = pd.DataFrame((rows), columns =['Type', 'Number'])
    
    # Creating pie plot from number of correct and incorrect attributes and saving it as png
    plot = statistic.groupby(['Type']).sum().plot(kind='pie', y='Number', autopct='%1.0f%%', label='').get_figure()
    plt.savefig(f"image1-{xmlfile[35:]}.png", bbox_inches='tight')
    
    # Grouping correct attributes by Area, saving as dataframe and csv file
    groupedc = matched.groupby(['Area']).count()
    groupedc = pd.DataFrame(groupedc)
    groupedc.to_csv(f'matched grouped-{xmlfile[35:]}.csv') 
    groupedc = pd.read_csv(f'matched grouped-{xmlfile[35:]}.csv')
    groupedc = groupedc[['Area', 'Comment']]
    
    # Grouping incorrect attributes by Area, saving as dataframe and csv file
    groupedi = log.groupby(['Area']).count()
    groupedi = pd.DataFrame(groupedi)
    groupedi.to_csv(f'incorrect grouped-{xmlfile[35:]}.csv') 
    groupedi = pd.read_csv(f'incorrect grouped-{xmlfile[35:]}.csv')
    groupedi = groupedi[['Area', 'Comment']]

    # Merging correct and incorrect attributes' dataframes and preparing for creating the table
    df_cd = pd.merge(groupedc, groupedi, how='outer', on='Area')
    df_cd['Comment_y']=df_cd['Comment_y'].fillna(0)
    df_cd = df_cd.rename(columns= {'Comment_x':'Number of matched attributes','Comment_y':'Number of mismatched attributes'})

    # Creating new column with sum and percantages of number of correct and incorrect attributes
    df_cd['Number of all atributes'] = df_cd['Number of matched attributes'] + df_cd['Number of mismatched attributes']
    df_cd['Number of matched attributes'] = pd.to_numeric(df_cd['Number of matched attributes'], downcast='float')
    df_cd['Percentage of matched attributes'] = df_cd['Number of matched attributes']*100 / df_cd['Number of all atributes']
    df_cd['Percentage of mismatched attributes'] = df_cd['Number of mismatched attributes']*100 / df_cd['Number of all atributes']

     # Creating barplot from percantages of number of correct and incorrect attributes- saving as png
    df_cd.plot(x="Area", y=["Percentage of matched attributes", "Percentage of mismatched attributes"], kind="bar")
    plt. ylabel("Percantage of attributes (%)")
    plt.savefig(f"image2-{xmlfile[35:]}.png", bbox_inches='tight')

    # Rounding values and adding '%'
    df_cd['Percentage of matched attributes'] = df_cd['Percentage of matched attributes'].round(2).astype(str) + "%"
    df_cd['Percentage of mismatched attributes'] = df_cd['Percentage of mismatched attributes'].round(2).astype(str) + "%"
    df_cd = df_cd.sort_values(by=['Number of all atributes'], ascending=False)
    

    # Exporting table to png
    dfi.export(df_cd,f"mytable2-{xmlfile[35:]}.png")

# PPTX EXPORT
# Opening presentation and creating first slide
ppt = Presentation('check results.pptx')
first_slide = ppt.slides.add_slide(ppt.slide_layouts[1])
title = "Check results  - " + str(date.today())

 # Setting the title of first slide and saving it in the presentation
first_slide.shapes.title.text = title
ppt.save('check results.pptx')

# Function for creating new slides and pasting tables and plots
def presentation(table1, table2, chart1, chart2, xmlfile):
    # Opening presentation
    ppt = Presentation('check results.pptx')

    # Creating titles for the slides
    title2 = f"Matched vs mismatched attributes in contract {xmlfile[35:]}"
    title3 = f"Mismatched attributes in contract {xmlfile[35:]}"
    title4 = f"Matched vs mismatched  attributes based on the area in contract {xmlfile[35:]}"
    title5 = f"Percentage of matched and mismatched attributes based on the area in contract {xmlfile[35:]}"
 
    # Setting the objects
    img = f'{chart1}.png'
    df = f'{table1}.png'
    img2 = f'{chart2}.png'
    df2 = f'{table2}.png'

    # Choosing the layout of the slide
    second_slide = ppt.slide_layouts[1]
    third_slide = ppt.slide_layouts[1]
    four_slide = ppt.slide_layouts[1]
    five_slide = ppt.slide_layouts[1]

    # Creating slides
    slide2 = ppt.slides.add_slide(second_slide)
    slide3 = ppt.slides.add_slide(third_slide)
    slide4 = ppt.slides.add_slide(four_slide)
    slide5 = ppt.slides.add_slide(five_slide)
    slide2.shapes.title.text = title2
    slide3.shapes.title.text = title3
    slide4.shapes.title.text = title4
    slide5.shapes.title.text = title5

    # Pasting charts and tables into the slides
    pic = slide2.shapes.add_picture(img, left= Inches(1),top = Inches(2),height = Inches(5))
    table = slide3.shapes.add_picture(df, left= Inches(2),top = Inches(3),height = Inches(2))
    table2 = slide4.shapes.add_picture(df2, left= Inches(1),top = Inches(3),height = Inches(2))
    pic2 = slide5.shapes.add_picture(img2, left= Inches(3),top = Inches(2),height = Inches(5))
        
    # Saving the powerpoint presentation
    ppt.save('check results.pptx')


# Loop for checking files and creating presentation
files = ['file1', 'file1']

for i in tqdm(range(len(files)), desc = colored("Update on the check process", "red"), colour="green"):
    check('file1', 
        'file2',
        files[i])
    presentation(f'mytable1-{files[i][35:]}', f'mytable2-{files[i][35:]}', f'image1-{files[i][35:]}', f'image2-{files[i][35:]}', files[i])
    sleep(0.1)
print(colored("Check done", "green"))
print(colored("Powerpoint presentation exported", "green"))

